"""Generate management PPT for TMS — narrative version (no cost/ROI).
Focus: Why it was built + How it benefits BCML.
Run: python docs/build_tms_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand colours ────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT_BLUE = RGBColor(0x2E, 0x6D, 0xA4)
GREEN       = RGBColor(0x05, 0x96, 0x69)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
RED         = RGBColor(0xDC, 0x26, 0x26)
GREY_LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
GREY_MID    = RGBColor(0x6B, 0x72, 0x80)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK   = RGBColor(0x0F, 0x17, 0x2A)


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, *,
              size=18, bold=False, color=TEXT_DARK,
              align=PP_ALIGN.LEFT, font='Calibri', anchor=MSO_ANCHOR.TOP):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tx


def _add_bullets(slide, left, top, width, height, lines, *,
                 size=14, color=TEXT_DARK, bullet_color=ACCENT_BLUE):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        b = p.add_run()
        b.text = "  •  "
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = 'Calibri'


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def _slide_header(slide, title, subtitle=None):
    _add_rect(slide, Emu(0), Emu(0), prs.slide_width, Inches(0.9), NAVY)
    _add_text(slide, Inches(0.5), Inches(0.18), Inches(12.5), Inches(0.55),
              title, size=24, bold=True, color=WHITE)
    if subtitle:
        _add_text(slide, Inches(0.5), Inches(1.05), Inches(12.5), Inches(0.45),
                  subtitle, size=14, color=GREY_MID)


def _footer(slide, n_of):
    _add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
              "TMS — Training Management System  ·  BCML", size=9, color=GREY_MID)
    _add_text(slide, Inches(10.5), Inches(7.05), Inches(3), Inches(0.3),
              f"{n_of}", size=9, color=GREY_MID, align=PP_ALIGN.RIGHT)


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ═══════════════ SLIDE 1 — Title ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
_set_bg(s, NAVY)
_add_rect(s, Emu(0), Inches(3.2), prs.slide_width, Inches(0.08), AMBER)
_add_text(s, Inches(0.5), Inches(2.0), Inches(12), Inches(0.5),
          "BCML — Balrampur Chini Mills Ltd.",
          size=14, color=AMBER, bold=True)
_add_text(s, Inches(0.5), Inches(2.5), Inches(12), Inches(1),
          "Training Management System",
          size=54, bold=True, color=WHITE)
_add_text(s, Inches(0.5), Inches(3.6), Inches(12), Inches(0.6),
          "Why we built it. How it transforms BCML L&D.",
          size=20, color=GREY_LIGHT)
_add_text(s, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
          "Management Review  ·  May 2026",
          size=14, color=GREY_LIGHT)


# ═══════════════ SLIDE 2 — Why we built this ═════════════════════════════
s = prs.slides.add_slide(blank)
_slide_header(s, "Why we built TMS",
              "Training records are the foundation of workforce capability — yet they lived in fragmented Excel files")

challenges = [
    ("Data scattered across plants",
     "Every plant maintained its own Excel sheets. Names, formats and definitions drifted apart. Group-level reporting required weeks of consolidation."),
    ("Free-text errors silently distorted everything",
     "Programme name spelt three ways became three different programmes. Coverage percentages reflected spelling, not learning."),
    ("Paper attendance, manual transcription",
     "Signed registers were re-typed into Excel each month, opening room for error and slow turnaround."),
    ("No single source of truth for audits",
     "ISO, customer and internal audits triggered week-long file-search exercises before each visit."),
    ("Decisions made on stale information",
     "Plant heads and Central L&D worked on month-old data. Course corrections were always retrospective."),
    ("Workforce-capability story disconnected from business strategy",
     "No way to link training to skill-gap closure, automation readiness or compliance posture."),
]
y = Inches(1.7)
for h, d in challenges:
    _add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.78), AMBER)
    _add_rect(s, Inches(0.65), y, Inches(12.2), Inches(0.78), GREY_LIGHT)
    _add_text(s, Inches(0.85), y + Inches(0.06), Inches(11.9), Inches(0.4),
              h, size=13, bold=True, color=NAVY)
    _add_text(s, Inches(0.85), y + Inches(0.4), Inches(11.9), Inches(0.5),
              d, size=11, color=GREY_MID)
    y += Inches(0.85)

_footer(s, "2 / 9")


# ═══════════════ SLIDE 3 — What TMS is ═════════════════════════════════════
s = prs.slides.add_slide(blank)
_slide_header(s, "TMS — a single platform across all 10 BCML plants",
              "Replaces fragmented Excel with one connected, real-time training system")

modules = [
    ("Employee Master",       "Live headcount with collar and role"),
    ("Yearly TNI",            "Single demand list with hygiene engine"),
    ("Programme Master",      "Canonical catalogue, no spelling drift"),
    ("Training Calendar",     "Plan, schedule, track, close"),
    ("Attendance (2A)",       "Validated entry, soon QR-based"),
    ("Programme Details (2C)","Gated saves, no garbage in"),
    ("Monthly Summary",       "Auto-built, real-time"),
    ("Central Dashboard",     "All plants, one view"),
    ("Compliance Engine",     "BC and WC coverage at a glance"),
    ("Audit Trail",           "Who did what, when, from where"),
    ("Excel Export",          "BCML's audit-format pack on demand"),
    ("Data Quality Engine",   "4-layer hygiene before any insert"),
    ("Public QR (Foundation)","Attendance, feedback, live view"),
]
cols = 3
col_w = Inches(4.1)
row_h = Inches(0.85)
start_left = Inches(0.4)
start_top  = Inches(1.7)
for i, (name, desc) in enumerate(modules):
    r = i // cols; c = i % cols
    x = start_left + c * (col_w + Inches(0.1))
    y = start_top + r * (row_h + Inches(0.1))
    _add_rect(s, x, y, col_w, row_h, GREY_LIGHT)
    _add_rect(s, x, y, Inches(0.06), row_h, GREEN)
    _add_text(s, x + Inches(0.2), y + Inches(0.08), col_w - Inches(0.3), Inches(0.4),
              name, size=12, bold=True, color=NAVY)
    _add_text(s, x + Inches(0.2), y + Inches(0.45), col_w - Inches(0.3), Inches(0.4),
              desc, size=10, color=GREY_MID)

_footer(s, "3 / 9")


# ═══════════════ SLIDE 4 — End-to-end Flow ═════════════════════════════════
s = prs.slides.add_slide(blank)
_slide_header(s, "End-to-end workflow",
              "From training need to compliance dashboard — one connected flow")

steps = ["Employees", "TNI", "Programme\nMaster", "Calendar", "2A", "2C", "Summary", "Dashboard", "Export"]
step_w = Inches(1.25); step_h = Inches(1.1); gap = Inches(0.12)
total_w = len(steps) * step_w + (len(steps) - 1) * gap
left0 = (prs.slide_width - total_w) // 2
top0  = Inches(2.5)
colors = [GREEN, ACCENT_BLUE, ACCENT_BLUE, ACCENT_BLUE, AMBER, AMBER, NAVY, NAVY, GREEN]
for i, step in enumerate(steps):
    x = left0 + i * (step_w + gap)
    _add_rect(s, x, top0, step_w, step_h, colors[i])
    tx = s.shapes.add_textbox(x, top0, step_w, step_h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = step
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
    if i < len(steps) - 1:
        ax = x + step_w + Emu(2000)
        ay = top0 + step_h / 2 - Inches(0.1)
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, gap - Emu(4000), Inches(0.2))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = GREY_MID
        arrow.line.fill.background()

_add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.4),
          "Each step reads from and writes to a shared database. Nothing is typed twice.",
          size=14, color=GREY_MID, align=PP_ALIGN.CENTER)

_add_rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.6), GREY_LIGHT)
_add_text(s, Inches(0.7), Inches(5.15), Inches(12), Inches(0.4),
          "Role-based access keeps data clean and visible to the right people",
          size=13, bold=True, color=NAVY)
_add_bullets(s, Inches(0.7), Inches(5.55), Inches(12), Inches(1),
             ["SPOC — own plant only; uploads, marks attendance, files 2C, views own summary",
              "Central L&D — all 10 plants in one view; approves back-fill requests, monitors anomalies",
              "Admin — full system; audit log; manages users and plant overrides"],
             size=12, bullet_color=NAVY)
_footer(s, "4 / 9")


# ═══════════════ SLIDE 5 — Data Quality Engine ═════════════════════════════
s = prs.slides.add_slide(blank)
_slide_header(s, "Data Quality Engine",
              "Why this matters — clean data is the foundation of every analytic, every decision")

layers = [
    ("1. Normalisation", "Whitespace, casing, unicode, curly quotes, trailing notes stripped", ACCENT_BLUE),
    ("2. Validation",    "Placeholder entries (NA, ?, TBD) rejected at the door",              NAVY),
    ("3. Spell-check",   "English dictionary correction, domain-aware allowlist",              GREEN),
    ("4. Master match",  "Each name reconciled with your Programme Master",                    AMBER),
]
ly = Inches(1.7)
for name, desc, col in layers:
    _add_rect(s, Inches(0.5), ly, Inches(0.15), Inches(0.95), col)
    _add_rect(s, Inches(0.65), ly, Inches(5.5), Inches(0.95), GREY_LIGHT)
    _add_text(s, Inches(0.85), ly + Inches(0.1), Inches(5.3), Inches(0.4),
              name, size=14, bold=True, color=NAVY)
    _add_text(s, Inches(0.85), ly + Inches(0.5), Inches(5.3), Inches(0.4),
              desc, size=11, color=GREY_MID)
    ly += Inches(1.1)

_add_text(s, Inches(6.7), Inches(1.7), Inches(6.3), Inches(0.4),
          "Real examples — caught before insert",
          size=14, bold=True, color=NAVY)

table_data = [
    ("BEFORE (typed by HoD)",          "AFTER (auto-fixed)"),
    ("Thept Prevention",                "Theft Prevention"),
    ("Saftey Procidure",                "Safety Procedure"),
    ("Maitenance of Equipement",        "Maintenance of Equipment"),
    ("Reportng of Non EHS Lapses",      "Reporting of Non EHS Lapses"),
    ("Lifting Tool and Takels",         "Lifting Tool and Tackles"),
    ("NA  /  ?  /  TBD",                "Rejected as placeholder"),
]
table = s.shapes.add_table(rows=len(table_data), cols=2,
                           left=Inches(6.7), top=Inches(2.2),
                           width=Inches(6.3), height=Inches(4.0)).table
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(3.3)
for ri, row in enumerate(table_data):
    for ci, val in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.color.rgb = WHITE if ri == 0 else TEXT_DARK
                r.font.bold = (ri == 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else (
            GREY_LIGHT if ri % 2 else WHITE)
_footer(s, "5 / 9")


# ═══════════════ SLIDE 6 — Anti-tamper roadmap ════════════════════════════
s = prs.slides.add_slide(blank)
_slide_header(s, "Anti-tamper roadmap",
              "Phase 2 — making every attendance record audit-grade")

defences = [
    ("QR self-scan",         "Employee scans on own phone at venue — eliminates paper register"),
    ("GPS lock",             "Scan rejected if not within registered venue range"),
    ("Time window",          "Scan only valid within session start and end window"),
    ("In-session codes",     "Trainer issues codes during session; missed codes reduce credit"),
    ("Faculty email confirm","External faculty signs off attendee list via secure one-time link"),
    ("Mentor sign-off (OJT)","For multi-day OJT, mentor confirms attendance daily via own login"),
    ("Hash chain audit",     "Every record cryptographically linked — silent edits are detectable"),
]
left_col = Inches(0.5)
top_band = Inches(1.7)
band_h   = Inches(0.65)
for i, (name, desc) in enumerate(defences):
    y = top_band + i * (band_h + Inches(0.08))
    _add_rect(s, left_col, y, Inches(0.45), band_h, GREEN)
    _add_text(s, left_col, y, Inches(0.45), band_h,
              str(i + 1), size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    _add_rect(s, left_col + Inches(0.45), y, Inches(12.4), band_h, GREY_LIGHT)
    _add_text(s, left_col + Inches(0.7), y + Inches(0.06), Inches(3.5), Inches(0.4),
              name, size=13, bold=True, color=NAVY)
    _add_text(s, left_col + Inches(4.3), y + Inches(0.06), Inches(8.4), Inches(0.55),
              desc, size=12, color=GREY_MID)
_footer(s, "6 / 9")


# ═══════════════ SLIDE 7 — How it benefits BCML ════════════════════════════
s = prs.slides.add_slide(blank)
_slide_header(s, "How TMS benefits BCML",
              "Six tangible shifts across the L&D function and beyond")

benefits = [
    ("Single source of truth",        "All training data in one place — plant, central, board all see the same numbers",         NAVY),
    ("Confidence in compliance",      "Coverage %, attendance and capability data accurate, defensible and audit-ready",          GREEN),
    ("Faster decisions",              "Real-time dashboards mean course-corrections happen mid-cycle, not after the damage",      ACCENT_BLUE),
    ("Empowered SPOCs",               "Hours of admin replaced by minutes — time freed for skill development, faculty management", AMBER),
    ("Workforce-capability visibility", "Skill gaps, programme effectiveness and reskilling velocity visible across the group",     ACCENT_BLUE),
    ("Industry 4.0 ready",            "Clean, machine-readable data — foundation for AI nominations, predictive planning, integration", GREEN),
]
cols = 2
cw = Inches(6.2)
ch = Inches(1.55)
sx = Inches(0.5)
sy = Inches(1.8)
for i, (title, desc, col) in enumerate(benefits):
    row = i // cols; c = i % cols
    x = sx + c * (cw + Inches(0.2))
    y = sy + row * (ch + Inches(0.18))
    _add_rect(s, x, y, cw, ch, GREY_LIGHT)
    _add_rect(s, x, y, Inches(0.15), ch, col)
    _add_text(s, x + Inches(0.35), y + Inches(0.18), cw - Inches(0.5), Inches(0.5),
              title, size=15, bold=True, color=NAVY)
    _add_text(s, x + Inches(0.35), y + Inches(0.68), cw - Inches(0.5), Inches(0.9),
              desc, size=11, color=GREY_MID)
_footer(s, "7 / 9")


# ═══════════════ SLIDE 8 — Roadmap ═════════════════════════════════════════
s = prs.slides.add_slide(blank)
_slide_header(s, "Where we go next",
              "Phased delivery; each sprint scoped and ready to start")

quarters = [
    ("Live today",     "All core modules in production. Smart Analyzer. Programme Master Duplicate Merger. Daily use across 10 plants.", GREEN),
    ("Sprint 1",       "Anti-tamper: QR check-in / out, GPS lock, time window. Faculty email confirm. OJT mentor sign-off.",  ACCENT_BLUE),
    ("Sprint 2",       "Audit hardening: hash chain, auto-lock after N days, edit-creates-request workflow.",                  NAVY),
    ("Sprint 3",       "Intelligence: predictive nominations, skill-gap analytics, HRMS integration, compliance forecast.",     AMBER),
    ("Sprint 4",       "Scale: enterprise database, mobile-first interface, multi-FY archival, advanced reports.",              GREY_MID),
]
qy = Inches(1.8)
for q, desc, col in quarters:
    _add_rect(s, Inches(0.5), qy, Inches(2.8), Inches(0.95), col)
    _add_text(s, Inches(0.5), qy + Inches(0.25), Inches(2.8), Inches(0.5),
              q, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_rect(s, Inches(3.4), qy, Inches(9.5), Inches(0.95), GREY_LIGHT)
    _add_text(s, Inches(3.6), qy + Inches(0.25), Inches(9.2), Inches(0.5),
              desc, size=12, color=TEXT_DARK)
    qy += Inches(1.05)
_footer(s, "8 / 9")


# ═══════════════ SLIDE 9 — Ask ═════════════════════════════════════════════
s = prs.slides.add_slide(blank)
_slide_header(s, "What we need from leadership",
              "Decisions to take TMS to the next level")

asks = [
    ("Endorsement of the journey",     "Recognise TMS as the canonical L&D system; deprecate parallel Excel-based tracking."),
    ("Approval for Sprint 1",          "Green-light anti-tamper QR + GPS rollout; pilot at one plant, then scale."),
    ("Plant SPOC ownership",           "Confirm SPOC at each plant with dedicated time blocks for TMS use."),
    ("HRMS integration",               "Authorise Darwinbox API access for employee master auto-sync."),
    ("Monthly review cadence",         "Approve a 30-min monthly Central review of adoption, coverage and anomalies."),
]
for i, (h, d) in enumerate(asks):
    y = Inches(1.8) + i * Inches(0.95)
    _add_rect(s, Inches(0.5), y, Inches(0.6), Inches(0.8), AMBER)
    _add_text(s, Inches(0.5), y, Inches(0.6), Inches(0.8),
              str(i + 1), size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    _add_rect(s, Inches(1.15), y, Inches(11.7), Inches(0.8), GREY_LIGHT)
    _add_text(s, Inches(1.35), y + Inches(0.08), Inches(11.3), Inches(0.4),
              h, size=14, bold=True, color=NAVY)
    _add_text(s, Inches(1.35), y + Inches(0.42), Inches(11.3), Inches(0.4),
              d, size=12, color=GREY_MID)

_add_rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6), NAVY)
_add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6),
          "TMS is live. With this support, BCML L&D becomes group-level, real-time and Industry 4.0 ready.",
          size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)
_footer(s, "9 / 9")


out_path = r"docs\TMS_Management_Review.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
